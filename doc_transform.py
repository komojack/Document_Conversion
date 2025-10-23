import os
import sys
import argparse
import json
from datetime import datetime

_has_docx = False
_has_pypdf2 = False
_has_pptx = False
_has_bs4 = False

try:
    from docx import Document
    _has_docx = True
except Exception:
    pass

try:
    from PyPDF2 import PdfReader
    _has_pypdf2 = True
except Exception:
    pass

try:
    from pptx import Presentation
    from pptx.util import Pt
    _has_pptx = True
except Exception:
    pass

try:
    from bs4 import BeautifulSoup
    _has_bs4 = True
except Exception:
    pass

# OCR 相关可选依赖
_has_pymupdf = False
_has_paddleocr = False
_has_numpy = False
_has_pil = False
_ocr_instance = None

try:
    import fitz  # PyMuPDF
    _has_pymupdf = True
except Exception:
    pass

try:
    from paddleocr import PaddleOCR
    _has_paddleocr = True
except Exception:
    pass

try:
    import numpy as np
    _has_numpy = True
except Exception:
    pass

try:
    from PIL import Image
    import io
    _has_pil = True
except Exception:
    pass

# PDF 生成相关可选依赖
_has_reportlab = False
_has_markdown = False

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    _has_reportlab = True
except Exception:
    pass

try:
    import markdown
    _has_markdown = True
except Exception:
    pass

import lazyllm

SUPPORTED_INPUTS = {".docx", ".pdf", ".pptx", ".md", ".txt", ".html", ".htm"}
SUPPORTED_OUTPUTS = {"md", "html", "pptx", "docx", "pdf"}


def _read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# ============== PDF OCR 辅助 ==============

def get_ocr(lang_preference: str = "ch"):
    global _ocr_instance
    if _ocr_instance is None:
        if not _has_paddleocr:
            raise RuntimeError("未安装 paddleocr，请先 `pip install paddleocr`。")
        try:
            _ocr_instance = PaddleOCR(use_angle_cls=True, lang=lang_preference, show_log=False)
        except Exception:
            # 回退英文模型
            _ocr_instance = PaddleOCR(use_angle_cls=True, lang="en", show_log=False)
    return _ocr_instance


def extract_text_from_pdf_ocr(path: str, dpi: int = 200, lang_preference: str = "ch") -> str:
    if not _has_pymupdf:
        raise RuntimeError("未安装 PyMuPDF，请先 `pip install pymupdf`（包名 PyMuPDF）。")
    if not (_has_numpy and _has_paddleocr and _has_pil):
        raise RuntimeError("OCR 所需依赖未就绪，请安装：numpy、paddleocr、Pillow。")

    ocr = get_ocr(lang_preference)
    doc = fitz.open(path)
    parts: list[str] = []

    for page in doc:
        pix = page.get_pixmap(dpi=dpi)
        try:
            img_bytes = pix.tobytes("png")
        except Exception:
            img_bytes = pix.getPNGData()
        # 使用PIL替代cv2处理图像
        try:
            pil_img = Image.open(io.BytesIO(img_bytes))
            # 转换为RGB格式（PaddleOCR需要）
            if pil_img.mode != 'RGB':
                pil_img = pil_img.convert('RGB')
            # 转换为numpy数组
            img = np.array(pil_img)
        except Exception:
            continue

        # PaddleOCR 返回结构：[[[box, (text, score)], ...]]
        try:
            result = ocr.ocr(img, cls=True)
        except Exception:
            result = []

        items = []
        for page_res in (result or []):
            for line in page_res:
                try:
                    box = line[0]
                    text = line[1][0]
                    # 以包围盒中心坐标排序，保证自上而下、从左到右
                    y = sum(pt[1] for pt in box) / len(box)
                    x = sum(pt[0] for pt in box) / len(box)
                    items.append((y, x, text))
                except Exception:
                    pass
        items.sort()  # 先按 y 再按 x
        lines = [t for (_, __, t) in items if t]
        if lines:
            parts.append("\n".join(lines))

    return "\n\n".join(parts)


# ============== 文本抽取 ==============

def extract_text_from_docx(path: str) -> str:
    if not _has_docx:
        raise RuntimeError("未安装 python-docx，请先 `pip install python-docx`。")
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        parts.append(p.text)
    return "\n".join(parts)


def extract_text_from_pdf(path: str) -> str:
    # 先尝试可复制文本的直读；若失败或为空，则回退 OCR
    text = ""
    if _has_pypdf2:
        reader = PdfReader(path)
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
        text = "\n".join(parts)

    if text.strip():
        return text

    # 走 OCR 回退路径（优先中文，失败回退英文）
    try:
        return extract_text_from_pdf_ocr(path, dpi=220, lang_preference="ch")
    except Exception:
        return extract_text_from_pdf_ocr(path, dpi=220, lang_preference="en")


def extract_text_from_pptx(path: str) -> str:
    if not _has_pptx:
        raise RuntimeError("未安装 python-pptx，请先 `pip install python-pptx`。")
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    tf = shape.text_frame
                    if tf and tf.text:
                        parts.append(tf.text)
            except Exception:
                # 忽略不可解析形状
                pass
    return "\n".join(parts)


def extract_text_from_html(path: str) -> str:
    if not _has_bs4:
        raise RuntimeError("未安装 beautifulsoup4，请先 `pip install beautifulsoup4`。")
    html = _read_text_file(path)
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n")


def extract_text_from_file(path: str) -> str:
    if not os.path.exists(path):
        raise FileNotFoundError(f"输入文件不存在: {path}")
    if os.path.isdir(path):
        raise IsADirectoryError(f"输入路径是目录而非文件: {path}")
    _, ext = os.path.splitext(path.lower())
    if ext not in SUPPORTED_INPUTS:
        raise ValueError(f"不支持的输入格式: {ext}, 支持: {sorted(SUPPORTED_INPUTS)}")

    if ext == ".txt" or ext == ".md":
        return _read_text_file(path)
    if ext == ".docx":
        return extract_text_from_docx(path)
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext == ".pptx":
        return extract_text_from_pptx(path)
    if ext in {".html", ".htm"}:
        return extract_text_from_html(path)

    # 兜底（理论不应到达）
    return _read_text_file(path)


def get_llm():

    return lazyllm.OnlineChatModule(source="sensenova", model="DeepSeek-V3-1")


def build_prompt(text: str, target_format: str, style: str | None = None) -> tuple[str, dict]:
    """返回 (instruction_prompt, extra_keys)
    pptx 目标将要求输出 JSON（slides）。docx 目标将先生成 Markdown。
    """
    style_desc = "风格：" + (style or "default")

    if target_format == "md":
        instruction = (
            "你是专业的文档编辑器。请将提供的原文重整为结构清晰的 Markdown，"
            "要求：\n"
            "- 顶级标题反映主题，合理分层（#、##、###）；\n"
            "- 保留重要要点为列表；\n"
            "- 对代码片段使用三引号代码块；\n"
            "- 表格以 Markdown 表格形式给出（若需要）；\n"
            "- 不要添加与原文无关的内容。\n"
            f"{style_desc}\n\n原文如下：\n{{context_str}}"
        )
        return instruction, {"context_str": text}

    if target_format == "html":
        instruction = (
            "你是专业的文档编辑器。请将原文转换为语义化 HTML（不包含多余样式），"
            "结构合理（h1/h2/h3、p、ul/ol、pre/code、table 等）。"
            "仅输出完整 HTML 文本，不要解释。\n"
            f"{style_desc}\n\n原文如下：\n{{context_str}}"
        )
        return instruction, {"context_str": text}

    if target_format == "pptx":
        instruction = (
            "你是资深演示文稿设计师。请将原文提炼为演示文稿提纲，"
            "仅输出 JSON，格式：{\"slides\": [{\"title\": \"...\", \"bullets\": [\"...\"]}]}.\n"
            "要求：\n- 5~10 页幻灯片；\n- 每页 3~6 条要点；\n- 不要包含除 JSON 外的任何文本。\n"
            f"{style_desc}\n\n原文如下：\n{{context_str}}"
        )
        return instruction, {"context_str": text}

    if target_format == "docx":
        # 先生成 Markdown，再渲染为 docx
        instruction = (
            "你是专业的文档编辑器。请将原文重整为结构清晰的 Markdown（后续将被转换为 DOCX），"
            "要求同 Markdown 目标，且只输出 Markdown 文本。\n"
            f"{style_desc}\n\n原文如下：\n{{context_str}}"
        )
        return instruction, {"context_str": text}

    if target_format == "pdf":
        # 先生成 Markdown，再渲染为 PDF
        instruction = (
            "你是专业的文档编辑器。请将原文重整为结构清晰的 Markdown（后续将被转换为 PDF），"
            "要求同 Markdown 目标，且只输出 Markdown 文本。\n"
            f"{style_desc}\n\n原文如下：\n{{context_str}}"
        )
        return instruction, {"context_str": text}

    raise ValueError(f"不支持的目标格式: {target_format}")


def llm_transform(text: str, target_format: str, style: str | None = None):
    llm = get_llm()
    instruction, extra = build_prompt(text, target_format, style)
    llm.prompt(lazyllm.ChatPrompter(instruction=instruction, extra_keys=list(extra.keys())))
    res = llm(extra)
    # 约定：返回字符串；若为 pptx 则需要解析为 JSON
    if target_format == "pptx":
        # 兼容 LLM 返回包含代码块围栏或额外说明的情况
        raw = res.strip()
        # 去除开头的 ``` 或 ```json 围栏
        if raw.startswith("```"):
            if raw.startswith("```json"):
                raw = raw[len("```json"):].lstrip("\n")
            else:
                raw = raw[3:].lstrip("\n")
            # 去除结尾的 ``` 围栏
            if raw.endswith("```"):
                raw = raw[:-3].rstrip()
        # 若仍有非 JSON 前缀/后缀，提取最外层花括号包裹的 JSON 片段
        if not raw.strip().startswith("{") or not raw.strip().endswith("}"):
            start = raw.find("{")
            end = raw.rfind("}")
            if start != -1 and end != -1 and end > start:
                raw = raw[start:end+1]
        try:
            data = json.loads(raw)
            if not isinstance(data, dict) or "slides" not in data:
                raise ValueError("JSON 结构不含 slides")
            return data
        except Exception as e:
            raise RuntimeError(f"LLM 未返回有效 JSON：{e}\n原始返回：{res[:500]}...")
    else:
        return res


# ============== 输出写入器 ==============

def write_markdown(text: str, output_path: str) -> str:
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    return output_path


def write_html(text: str, output_path: str) -> str:
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    return output_path


def write_pptx(slides_json: dict, output_path: str) -> str:
    if not _has_pptx:
        raise RuntimeError("未安装 python-pptx，请先 `pip install python-pptx`。")
    prs = Presentation()
    # 使用标题+内容布局（一般是索引 1），但不同模板可能不同；使用第一个内容布局兜底
    layout = None
    try:
        layout = prs.slide_layouts[1]
    except Exception:
        layout = prs.slide_layouts[0]

    for slide_spec in slides_json.get("slides", []):
        title = slide_spec.get("title", "")
        bullets = slide_spec.get("bullets", [])
        slide = prs.slides.add_slide(layout)
        try:
            slide.shapes.title.text = title
        except Exception:
            pass
        # 内容框（一般第二个占位）
        body_shape = None
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
        # 若不可用，则尝试第一个形状文本框
        text_frame = None
        if body_shape and hasattr(body_shape, "text_frame"):
            text_frame = body_shape.text_frame
            text_frame.clear()
        else:
            # 找到第一个支持 text_frame 的形状
            for shp in slide.shapes:
                if hasattr(shp, "text_frame"):
                    text_frame = shp.text_frame
                    text_frame.clear()
                    break
        if text_frame:
            # 第一行作为第一条要点
            for i, bullet in enumerate(bullets):
                p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                p.text = bullet
                p.level = 0
    prs.save(output_path)
    return output_path


def write_docx_from_markdown(md_text: str, output_path: str) -> str:
    if not _has_docx:
        raise RuntimeError("未安装 python-docx，请先 `pip install python-docx`。")
    doc = Document()
    in_code_block = False
    code_buffer = []

    def flush_code_block():
        nonlocal code_buffer
        if code_buffer:
            # 简单插入为普通段落（可扩展为设置等宽字体）
            for line in code_buffer:
                p = doc.add_paragraph()
                run = p.add_run(line)
                try:
                    run.font.name = "Courier New"
                except Exception:
                    pass
            code_buffer = []

    for raw_line in md_text.splitlines():
        line = raw_line.rstrip("\n")
        # 代码块处理
        if line.strip().startswith("```"):
            if not in_code_block:
                in_code_block = True
                code_buffer = []
            else:
                in_code_block = False
                flush_code_block()
            continue
        if in_code_block:
            code_buffer.append(line)
            continue

        # 标题处理
        if line.startswith("#"):
            # 计算#的数量作为标题级别
            level = 0
            for char in line:
                if char == '#':
                    level += 1
                else:
                    break
            title = line.lstrip('#').strip()
            try:
                # docx标题级别从0开始，最大为8，但通常使用0-3
                doc.add_heading(title, level=min(level, 4))
            except Exception:
                doc.add_paragraph(title)
            continue
        # 列表处理
        if line.strip().startswith("- ") or line.strip().startswith("* "):
            text = line.strip()[2:].strip()
            # 计算缩进级别
            indent_level = (len(line) - len(line.lstrip())) // 2
            p = doc.add_paragraph(text, style='List Bullet')
            # 设置缩进（如果有的话）
            if indent_level > 0:
                try:
                    p.paragraph_format.left_indent = indent_level * 720  # 720 twips = 0.5 inch
                except Exception:
                    pass
            continue
        # 空行
        if not line.strip():
            doc.add_paragraph("")
            continue
        # 普通段落
        doc.add_paragraph(line)

    # 尾部若仍在代码块，刷新
    if in_code_block:
        flush_code_block()

    doc.save(output_path)
    return output_path


def write_pdf(text: str, output_path: str, source_format: str = "markdown") -> str:
    """
    将文本转换为PDF格式
    
    Args:
        text: 输入文本内容
        output_path: 输出PDF文件路径
        source_format: 源格式，支持 "markdown" 或 "html"
    
    Returns:
        输出文件路径
    """
    if not _has_reportlab:
        raise RuntimeError("未安装 reportlab，请先 `pip install reportlab`。")
    
    # 注册中文字体
    def register_chinese_fonts():
        """注册中文字体，支持多种常见字体"""
        import platform
        system = platform.system()
        
        # 常见中文字体路径
        font_paths = []
        if system == "Windows":
            font_paths = [
                "C:/Windows/Fonts/simsun.ttc",  # 宋体
                "C:/Windows/Fonts/simhei.ttf",  # 黑体
                "C:/Windows/Fonts/msyh.ttc",    # 微软雅黑
                "C:/Windows/Fonts/simkai.ttf",  # 楷体
            ]
        elif system == "Darwin":  # macOS
            font_paths = [
                "/System/Library/Fonts/PingFang.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
                "/System/Library/Fonts/Hiragino Sans GB.ttc",
            ]
        else:  # Linux
            font_paths = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            ]
        
        # 尝试注册第一个可用的字体
        for font_path in font_paths:
            if os.path.exists(font_path):
                try:
                    pdfmetrics.registerFont(TTFont('ChineseFont', font_path))
                    return 'ChineseFont'
                except Exception:
                    continue
        
        # 如果没有找到系统字体，使用ReportLab内置的支持Unicode的字体
        try:
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
            return 'STSong-Light'
        except Exception:
            pass
        
        # 最后的回退选项
        return 'Helvetica'
    
    # 注册并获取中文字体名称
    chinese_font = register_chinese_fonts()
    
    # 创建PDF文档
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    
    # 自定义样式 - 使用支持中文的字体
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=18,
        spaceAfter=12,
        textColor=colors.HexColor('#2c3e50'),
        fontName=chinese_font
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=14,
        spaceAfter=8,
        textColor=colors.HexColor('#2c3e50'),
        fontName=chinese_font
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=10,
        spaceAfter=6,
        leading=14,
        fontName=chinese_font
    )
    
    code_style = ParagraphStyle(
        'CustomCode',
        parent=styles['Code'],
        fontSize=9,
        fontName='Courier',  # 代码保持等宽字体
        backColor=colors.HexColor('#f8f9fa'),
        borderColor=colors.HexColor('#3498db'),
        borderWidth=1,
        leftIndent=12,
        rightIndent=12,
        spaceAfter=6
    )
    
    if source_format == "markdown":
        if not _has_markdown:
            raise RuntimeError("未安装 markdown，请先 `pip install markdown`。")
        
        # 将Markdown转换为HTML，然后解析为PDF内容
        html_content = markdown.markdown(
            text, 
            extensions=['tables', 'fenced_code', 'toc'],
        )
        
        # 简单的HTML解析和转换
        import re
        
        # 移除HTML标签并提取内容
        lines = html_content.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # 处理标题
            if line.startswith('<h1>') and line.endswith('</h1>'):
                title_text = re.sub(r'<[^>]+>', '', line)
                story.append(Paragraph(title_text, title_style))
                story.append(Spacer(1, 12))
            elif line.startswith('<h2>') and line.endswith('</h2>'):
                heading_text = re.sub(r'<[^>]+>', '', line)
                story.append(Paragraph(heading_text, heading_style))
                story.append(Spacer(1, 8))
            elif line.startswith('<h3>') and line.endswith('</h3>'):
                heading_text = re.sub(r'<[^>]+>', '', line)
                story.append(Paragraph(heading_text, heading_style))
                story.append(Spacer(1, 6))
            # 处理代码块
            elif line.startswith('<pre><code>') and line.endswith('</code></pre>'):
                code_text = re.sub(r'<[^>]+>', '', line)
                story.append(Paragraph(code_text, code_style))
                story.append(Spacer(1, 6))
            # 处理段落
            elif line.startswith('<p>') and line.endswith('</p>'):
                para_text = re.sub(r'<[^>]+>', '', line)
                if para_text:
                    story.append(Paragraph(para_text, normal_style))
                    story.append(Spacer(1, 6))
            # 处理列表项
            elif line.startswith('<li>') and line.endswith('</li>'):
                list_text = re.sub(r'<[^>]+>', '', line)
                story.append(Paragraph(f"• {list_text}", normal_style))
            # 处理其他内容
            elif not line.startswith('<') and line:
                story.append(Paragraph(line, normal_style))
                story.append(Spacer(1, 6))
                
    elif source_format == "html":
        # 简单的HTML处理
        import re
        clean_text = re.sub(r'<[^>]+>', '', text)
        lines = clean_text.split('\n')
        for line in lines:
            line = line.strip()
            if line:
                story.append(Paragraph(line, normal_style))
                story.append(Spacer(1, 6))
    else:
        # 直接处理纯文本
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if line:
                # 简单的Markdown标题检测
                if line.startswith('# '):
                    story.append(Paragraph(line[2:], title_style))
                    story.append(Spacer(1, 12))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], heading_style))
                    story.append(Spacer(1, 8))
                elif line.startswith('### '):
                    story.append(Paragraph(line[4:], heading_style))
                    story.append(Spacer(1, 6))
                else:
                    story.append(Paragraph(line, normal_style))
                    story.append(Spacer(1, 6))
    
    # 如果没有内容，添加默认内容
    if not story:
        story.append(Paragraph("文档内容", normal_style))
    
    # 生成PDF
    try:
        doc.build(story)
    except Exception as e:
        raise RuntimeError(f"PDF生成失败: {e}")
    
    return output_path


# ============== CLI ==============

def infer_output_path(input_path: str, target_format: str, user_output: str | None) -> str:
    if user_output:
        return os.path.abspath(user_output)
    base = os.path.splitext(os.path.basename(input_path))[0]
    dir_ = os.path.dirname(os.path.abspath(input_path))
    suffix = {"md": ".md", "html": ".html", "pptx": ".pptx", "docx": ".docx", "pdf": ".pdf"}[target_format]
    return os.path.join(dir_, f"{base}_converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}{suffix}")


def main():
    parser = argparse.ArgumentParser(description="基于AI的文档格式转换工具")
    parser.add_argument("--input", required=True, help="输入文件路径（docx/pdf/pptx/md/txt/html）")
    parser.add_argument("--to", required=True, choices=sorted(SUPPORTED_OUTPUTS), help="目标格式：md/html/pptx/docx/pdf")
    parser.add_argument("--output", required=False, help="输出文件路径（默认与输入同目录自动命名）")
    parser.add_argument("--style", required=False, help="转换风格（如 concise/academic 等）")

    args = parser.parse_args()

    input_path = os.path.abspath(args.input)
    target = args.to
    output_path = infer_output_path(input_path, target, args.output)

    print({"input": input_path, "to": target, "output": output_path, "style": args.style})

    try:
        print("[1/4] 抽取原文文本...")
        text = extract_text_from_file(input_path)
        if not text.strip():
            print("警告：抽取到的文本为空，可能为扫描版 PDF 或复杂格式。")
        print("[2/4] 触发 LLM 转换...")
        transformed = llm_transform(text, target, args.style)
        print("[3/4] 写出结果...")
        if target == "md":
            write_markdown(transformed, output_path)
        elif target == "html":
            write_html(transformed, output_path)
        elif target == "pptx":
            write_pptx(transformed, output_path)
        elif target == "docx":
            write_docx_from_markdown(transformed, output_path)
        elif target == "pdf":
            write_pdf(transformed, output_path, source_format="markdown")
        else:
            raise ValueError(f"未知目标格式: {target}")
        print(f"[4/4] 完成：{output_path}")
    except Exception as e:
        print(f"错误：{e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()